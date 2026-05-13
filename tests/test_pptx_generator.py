"""Tests for src/pptx_generator.py PPTX report generator."""

import glob
import os
import tempfile

import pandas as pd
import pytest
from pptx import Presentation

from src.ig_processor import process as ig_process
from src.pptx_generator import generate_report
from src.x_processor import process as x_process

SAMPLE_IG_CSV = os.path.join(
    os.path.dirname(__file__), "sample_data", "sample_ig_content.csv"
)
SAMPLE_X_CSV = os.path.join(
    os.path.dirname(__file__), "sample_data", "sample_x_tweets.csv"
)


@pytest.fixture
def ig_df():
    return ig_process(SAMPLE_IG_CSV, "TestBrand")


@pytest.fixture
def x_df():
    return x_process(SAMPLE_X_CSV, "TestBrand")


@pytest.fixture
def output_dir():
    d = tempfile.mkdtemp()
    yield d
    # Cleanup any files created
    for f in glob.glob(os.path.join(d, "*")):
        os.remove(f)
    os.rmdir(d)


class TestGenerateReportIgOnly:
    def test_creates_pptx(self, ig_df, output_dir):
        path = os.path.join(output_dir, "report.pptx")
        result = generate_report(ig_df, None, "TestBrand", path)
        assert os.path.exists(result)

    def test_has_6_slides(self, ig_df, output_dir):
        path = os.path.join(output_dir, "report.pptx")
        generate_report(ig_df, None, "TestBrand", path)
        prs = Presentation(path)
        # Cover, Summary, IG Overview, IG Top Posts, Insights, Closing
        assert len(prs.slides) == 6


class TestGenerateReportBothPlatforms:
    def test_has_8_slides(self, ig_df, x_df, output_dir):
        path = os.path.join(output_dir, "report.pptx")
        generate_report(ig_df, x_df, "TestBrand", path)
        prs = Presentation(path)
        # Cover, Summary, IG Overview, IG Top Posts, X Overview, X Top Posts, Insights, Closing
        assert len(prs.slides) == 8

    def test_creates_pptx(self, ig_df, x_df, output_dir):
        path = os.path.join(output_dir, "report.pptx")
        result = generate_report(ig_df, x_df, "TestBrand", path)
        assert os.path.exists(result)


class TestGenerateReportXOnly:
    def test_has_6_slides(self, x_df, output_dir):
        path = os.path.join(output_dir, "report.pptx")
        generate_report(None, x_df, "TestBrand", path)
        prs = Presentation(path)
        # Cover, Summary, X Overview, X Top Posts, Insights, Closing
        assert len(prs.slides) == 6


class TestScorecardsCorrect:
    def test_summary_values_ig_only(self, ig_df, output_dir):
        path = os.path.join(output_dir, "report.pptx")
        generate_report(ig_df, None, "TestBrand", path)
        prs = Presentation(path)
        # Slide 2 is executive summary (index 1)
        summary_slide = prs.slides[1]
        all_text = " ".join(
            shape.text for shape in summary_slide.shapes if hasattr(shape, "text")
        )
        # Total posts should be 20
        assert "20" in all_text
        # Total reach should match sum
        total_reach = int(ig_df["reach"].sum())
        assert f"{total_reach:,}" in all_text

    def test_summary_values_both(self, ig_df, x_df, output_dir):
        path = os.path.join(output_dir, "report.pptx")
        generate_report(ig_df, x_df, "TestBrand", path)
        prs = Presentation(path)
        summary_slide = prs.slides[1]
        all_text = " ".join(
            shape.text for shape in summary_slide.shapes if hasattr(shape, "text")
        )
        # Total posts = 20 IG + 15 X = 35
        assert "35" in all_text


class TestTopPostsTable:
    def test_ig_top_5_sorted(self, ig_df, output_dir):
        path = os.path.join(output_dir, "report.pptx")
        generate_report(ig_df, None, "TestBrand", path)
        prs = Presentation(path)
        # IG top posts is slide 4 (index 3)
        top_posts_slide = prs.slides[3]

        # Find the table shape
        table_shape = None
        for shape in top_posts_slide.shapes:
            if shape.has_table:
                table_shape = shape
                break
        assert table_shape is not None

        # Extract engagement rates from the last column (Eng. Rate)
        table = table_shape.table
        rates = []
        for row_idx in range(1, table.rows.__len__()):
            rate_text = table.cell(row_idx, 4).text.rstrip("%")
            rates.append(float(rate_text))

        # Should be sorted descending
        assert rates == sorted(rates, reverse=True)

        # Should have at most 5 rows
        assert len(rates) <= 5


class TestInsightsGenerated:
    def test_contains_engagement_info(self, ig_df, output_dir):
        path = os.path.join(output_dir, "report.pptx")
        generate_report(ig_df, None, "TestBrand", path)
        prs = Presentation(path)
        # Insights slide: index 4 for IG-only (Cover, Summary, IG Overview, IG Top, Insights, Closing)
        insights_slide = prs.slides[4]
        all_text = " ".join(
            shape.text for shape in insights_slide.shapes if hasattr(shape, "text")
        )
        # Should mention engagement rate
        assert "engagement" in all_text.lower()
        # Should mention reach
        assert "reach" in all_text.lower()

    def test_contains_x_info_when_provided(self, ig_df, x_df, output_dir):
        path = os.path.join(output_dir, "report.pptx")
        generate_report(ig_df, x_df, "TestBrand", path)
        prs = Presentation(path)
        # Insights is slide 7 (index 6) with both platforms
        insights_slide = prs.slides[6]
        all_text = " ".join(
            shape.text for shape in insights_slide.shapes if hasattr(shape, "text")
        )
        assert "impressions" in all_text.lower()


class TestOutputFileExists:
    def test_returns_path(self, ig_df, output_dir):
        path = os.path.join(output_dir, "test_output.pptx")
        result = generate_report(ig_df, None, "TestBrand", path)
        assert result == path
        assert os.path.exists(path)
        # File should have content (> 0 bytes)
        assert os.path.getsize(path) > 0


class TestChartCleanup:
    def test_no_leftover_pngs(self, ig_df, x_df, output_dir):
        path = os.path.join(output_dir, "report.pptx")
        generate_report(ig_df, x_df, "TestBrand", path)
        # Check for any leftover chart PNGs in temp directories
        # The chart_dir should be cleaned up
        temp_dirs = glob.glob(os.path.join(tempfile.gettempdir(), "reportengine_charts_*"))
        for d in temp_dirs:
            pngs = glob.glob(os.path.join(d, "*.png"))
            assert len(pngs) == 0, f"Leftover PNGs found in {d}: {pngs}"
