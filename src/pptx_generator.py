"""PPTX report generator for ReportEngine.

Generates branded client-facing PowerPoint reports from processed
Instagram and X/Twitter DataFrames. Uses python-pptx for slide
construction and matplotlib for chart rendering.
"""

import logging
import os
import tempfile
from pathlib import Path
from typing import List, Optional, Tuple

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

logger = logging.getLogger(__name__)

# TNT brand colors
BLUE = RGBColor(30, 60, 115)
BEIGE = RGBColor(245, 240, 230)
DARK_TEXT = RGBColor(51, 51, 51)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(153, 153, 153)
LIGHT_BEIGE = RGBColor(250, 247, 240)

BLUE_HEX = "#1E3C73"
BEIGE_HEX = "#F5F0E6"

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Font
FONT_NAME = "Calibri"


def _set_slide_bg(slide, color: RGBColor) -> None:
    """Set the background color of a slide.

    Args:
        slide: python-pptx slide object.
        color: RGBColor for the background fill.
    """
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = DARK_TEXT,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """Add a text box to a slide.

    Args:
        slide: python-pptx slide object.
        left: Left position in EMU or Inches.
        top: Top position.
        width: Width.
        height: Height.
        text: Text content.
        font_size: Font size in points.
        bold: Whether text is bold.
        color: Text color.
        alignment: Text alignment.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.alignment = alignment


def _add_scorecard(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    value: str,
    label: str,
) -> None:
    """Add a scorecard box (rounded rectangle with value and label).

    Args:
        slide: python-pptx slide object.
        left: Left position.
        top: Top position.
        width: Width.
        height: Height.
        value: Large number/text to display.
        label: Small label below the value.
    """
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BLUE
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Value
    p_val = tf.paragraphs[0]
    p_val.text = value
    p_val.font.size = Pt(28)
    p_val.font.bold = True
    p_val.font.color.rgb = BLUE
    p_val.font.name = FONT_NAME
    p_val.alignment = PP_ALIGN.CENTER

    # Label
    p_label = tf.add_paragraph()
    p_label.text = label
    p_label.font.size = Pt(12)
    p_label.font.color.rgb = GRAY
    p_label.font.name = FONT_NAME
    p_label.alignment = PP_ALIGN.CENTER


def _create_bar_chart(
    data: dict,
    title: str,
    color: str,
    output_path: str,
    xlabel: str = "",
) -> str:
    """Create a horizontal bar chart with matplotlib and save as PNG.

    Args:
        data: Dict mapping labels to values.
        title: Chart title.
        color: Bar color as hex string.
        output_path: Path to save the PNG file.
        xlabel: X-axis label.

    Returns:
        The output_path where the PNG was saved.
    """
    fig, ax = plt.subplots(figsize=(5.5, 3.5))

    labels = list(data.keys())
    values = list(data.values())

    bars = ax.barh(labels, values, color=color, height=0.5)
    ax.set_title(title, fontsize=13, fontweight="bold", color="#1E3C73", pad=12,
                 fontfamily="sans-serif")

    if xlabel:
        ax.set_xlabel(xlabel, fontsize=10, fontfamily="sans-serif")

    ax.invert_yaxis()
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_visible(False)
    ax.tick_params(left=False, labelsize=10)
    ax.xaxis.grid(False)
    ax.yaxis.grid(False)

    for bar, val in zip(bars, values):
        label = f"{val:,.0f}" if val >= 100 else f"{val:.1f}%"
        ax.text(
            bar.get_width() + max(values) * 0.02,
            bar.get_y() + bar.get_height() / 2,
            label,
            va="center",
            fontsize=10,
            color="#333333",
            fontfamily="sans-serif",
        )

    ax.set_xlim(0, max(values) * 1.2 if values else 1)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    plt.tight_layout()
    plt.savefig(output_path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)

    return output_path


def _format_date_range(ig_df: Optional[pd.DataFrame], x_df: Optional[pd.DataFrame]) -> str:
    """Get human-readable date range from DataFrames.

    Args:
        ig_df: Instagram DataFrame or None.
        x_df: X DataFrame or None.

    Returns:
        Formatted date range string like "Mar 1 - Mar 28, 2026".
    """
    all_dates = []
    if ig_df is not None and len(ig_df) > 0:
        all_dates.extend(ig_df["date"].tolist())
    if x_df is not None and len(x_df) > 0:
        all_dates.extend(x_df["date"].tolist())

    if not all_dates:
        return "No data"

    dates = pd.to_datetime(all_dates)
    min_d = dates.min()
    max_d = dates.max()
    if min_d.year == max_d.year:
        return f"{min_d.strftime('%b %-d')} - {max_d.strftime('%b %-d, %Y')}"
    return f"{min_d.strftime('%b %-d, %Y')} - {max_d.strftime('%b %-d, %Y')}"


def _build_cover_slide(prs: Presentation, brand_name: str, date_range: str, period_label: str = "Monthly") -> None:
    """Build Slide 1: Cover."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    _set_slide_bg(slide, BEIGE)

    _add_textbox(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(1.2),
                 brand_name, font_size=36, bold=True, color=BLUE,
                 alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(0.8),
                 f"{period_label} Social Media Performance Report",
                 font_size=20, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.6),
                 date_range, font_size=16, color=DARK_TEXT,
                 alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.5),
                 "Prepared by TNT | Powered by Ascnd",
                 font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)


def _build_summary_slide(
    prs: Presentation,
    ig_df: Optional[pd.DataFrame],
    x_df: Optional[pd.DataFrame],
    deltas: Optional[dict] = None,
    goals_actuals: Optional[dict] = None,
) -> None:
    """Build Slide 2: Executive Summary with scorecards, deltas, and goals."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.7),
                 "Executive Summary", font_size=28, bold=True, color=BLUE)

    total_posts = 0
    total_reach = 0
    weighted_eng_sum = 0
    weighted_eng_denom = 0

    if ig_df is not None and len(ig_df) > 0:
        total_posts += len(ig_df)
        total_reach += int(ig_df["reach"].sum())
        weighted_eng_sum += (ig_df["engagement_rate"] * ig_df["reach"]).sum()
        weighted_eng_denom += ig_df["reach"].sum()

    if x_df is not None and len(x_df) > 0:
        total_posts += len(x_df)
        total_reach += int(x_df["impressions"].sum())
        weighted_eng_sum += (x_df["engagement_rate"] * x_df["impressions"]).sum()
        weighted_eng_denom += x_df["impressions"].sum()

    avg_eng = round(weighted_eng_sum / weighted_eng_denom, 2) if weighted_eng_denom > 0 else 0

    # Top content type by avg engagement rate
    top_type = "N/A"
    all_types = {}
    if ig_df is not None and len(ig_df) > 0:
        for ct, group in ig_df.groupby("content_type"):
            all_types[ct] = group["engagement_rate"].mean()
    if x_df is not None and len(x_df) > 0:
        for ct, group in x_df.groupby("content_type"):
            if ct not in all_types or group["engagement_rate"].mean() > all_types[ct]:
                all_types[ct] = group["engagement_rate"].mean()
    if all_types:
        top_type = max(all_types, key=all_types.get)

    card_width = Inches(2.7)
    card_height = Inches(1.8)
    card_top = Inches(2.2)
    gap = Inches(0.35)
    start_left = Inches(0.7)

    # Build delta suffix for reach and engagement if available
    def _delta_suffix(key: str) -> str:
        if not deltas or key not in deltas:
            return ""
        d = deltas[key]
        if d.get("delta_pct") is None:
            return ""
        pct = d["delta_pct"]
        arrow = "+" if pct > 0 else ""
        return f"\n{arrow}{pct}% vs prior"

    reach_delta = _delta_suffix("ig_reach")
    eng_delta = _delta_suffix("ig_engagement_rate")

    cards = [
        (f"{total_posts}", "Total Posts"),
        (f"{total_reach:,}{reach_delta}", "Total Reach"),
        (f"{avg_eng}%{eng_delta}", "Avg Engagement Rate"),
        (top_type, "Top Content Type"),
    ]

    for i, (value, label) in enumerate(cards):
        left = start_left + i * (card_width + gap)
        _add_scorecard(slide, left, card_top, card_width, card_height, value, label)

    # Goals vs Actuals row (if provided)
    if goals_actuals:
        _add_textbox(slide, Inches(0.7), Inches(4.4), Inches(11), Inches(0.4),
                     "Goals vs Actuals", font_size=14, bold=True, color=BLUE)
        goals_list = list(goals_actuals.items())
        g_card_w = Inches(2.7)
        g_card_h = Inches(1.2)
        for i, (g_label, g_data) in enumerate(goals_list[:4]):
            left = start_left + i * (card_width + gap)
            pct = g_data["pct"]
            met_str = "Met" if g_data["met"] else "Behind"
            val_str = f"{pct}% ({met_str})"
            _add_scorecard(slide, left, Inches(4.9), g_card_w, g_card_h, val_str, g_label)


def _build_ig_overview_slide(prs: Presentation, ig_df: pd.DataFrame, chart_dir: str) -> List[str]:
    """Build Slide 3: Instagram Overview with charts.

    Args:
        prs: Presentation object.
        ig_df: Instagram DataFrame.
        chart_dir: Temp directory for chart PNGs.

    Returns:
        List of chart PNG file paths created.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_files = []

    _add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.7),
                 "Instagram Performance", font_size=28, bold=True, color=BLUE)

    # Scorecards
    card_width = Inches(3.5)
    card_height = Inches(1.2)
    card_top = Inches(1.3)
    start_left = Inches(0.7)
    gap = Inches(0.4)

    ig_cards = [
        (f"{len(ig_df)}", "Posts"),
        (f"{int(ig_df['reach'].sum()):,}", "Total Reach"),
        (f"{round(ig_df['engagement_rate'].mean(), 2)}%", "Avg Engagement Rate"),
    ]
    for i, (value, label) in enumerate(ig_cards):
        left = start_left + i * (card_width + gap)
        _add_scorecard(slide, left, card_top, card_width, card_height, value, label)

    # Charts
    reach_by_type = ig_df.groupby("content_type")["reach"].mean().sort_values(ascending=True).to_dict()
    eng_by_type = ig_df.groupby("content_type")["engagement_rate"].mean().sort_values(ascending=True).to_dict()

    reach_chart = os.path.join(chart_dir, "ig_reach_by_type.png")
    _create_bar_chart(reach_by_type, "Avg Reach by Content Type", BLUE_HEX, reach_chart)
    chart_files.append(reach_chart)

    eng_chart = os.path.join(chart_dir, "ig_eng_by_type.png")
    _create_bar_chart(eng_by_type, "Avg Engagement Rate by Content Type", BLUE_HEX, eng_chart,
                      xlabel="")
    chart_files.append(eng_chart)

    slide.shapes.add_picture(reach_chart, Inches(0.5), Inches(3.2), Inches(6), Inches(4))
    slide.shapes.add_picture(eng_chart, Inches(6.8), Inches(3.2), Inches(6), Inches(4))

    return chart_files


def _build_top_posts_slide(
    slide,
    df: pd.DataFrame,
    columns: List[Tuple[str, str, float]],
) -> None:
    """Build a top posts table on an existing slide.

    Args:
        slide: python-pptx slide object.
        df: DataFrame sorted by engagement_rate descending.
        columns: List of (df_column, header_label, col_width_inches) tuples.
    """
    top5 = df.nlargest(5, "engagement_rate").reset_index(drop=True)

    rows = min(len(top5), 5) + 1  # +1 for header
    cols = len(columns)

    col_widths = [Inches(c[2]) for c in columns]
    table_width = sum(c[2] for c in columns)
    left = Inches((13.333 - table_width) / 2)

    table_shape = slide.shapes.add_table(rows, cols, left, Inches(1.6),
                                          Inches(table_width), Inches(0.5 * rows))
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    for i, (_, header, _) in enumerate(columns):
        cell = table.cell(0, i)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for row_idx in range(len(top5)):
        row_data = top5.iloc[row_idx]
        for col_idx, (col_name, _, _) in enumerate(columns):
            cell = table.cell(row_idx + 1, col_idx)
            val = row_data[col_name]
            if col_name == "engagement_rate":
                cell.text = f"{val}%"
            elif col_name in ("reach", "impressions"):
                cell.text = f"{int(val):,}"
            else:
                cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_TEXT
            p.font.name = FONT_NAME
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Alternating row colors
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BEIGE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE


def _build_ig_top_posts_slide(prs: Presentation, ig_df: pd.DataFrame) -> None:
    """Build Slide 4: Instagram Top Posts table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.7),
                 "Top Performing Instagram Posts", font_size=28, bold=True, color=BLUE)

    columns = [
        ("date", "Date", 1.8),
        ("content_type", "Type", 1.5),
        ("caption_preview", "Caption", 4.5),
        ("reach", "Reach", 1.8),
        ("engagement_rate", "Eng. Rate", 1.8),
    ]
    _build_top_posts_slide(slide, ig_df, columns)


def _build_x_overview_slide(prs: Presentation, x_df: pd.DataFrame, chart_dir: str) -> List[str]:
    """Build Slide 5: X Overview with charts.

    Args:
        prs: Presentation object.
        x_df: X DataFrame.
        chart_dir: Temp directory for chart PNGs.

    Returns:
        List of chart PNG file paths created.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_files = []

    _add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.7),
                 "X (Twitter) Performance", font_size=28, bold=True, color=BLUE)

    # Scorecards
    card_width = Inches(3.5)
    card_height = Inches(1.2)
    card_top = Inches(1.3)
    start_left = Inches(0.7)
    gap = Inches(0.4)

    x_cards = [
        (f"{len(x_df)}", "Tweets"),
        (f"{int(x_df['impressions'].sum()):,}", "Total Impressions"),
        (f"{round(x_df['engagement_rate'].mean(), 2)}%", "Avg Engagement Rate"),
    ]
    for i, (value, label) in enumerate(x_cards):
        left = start_left + i * (card_width + gap)
        _add_scorecard(slide, left, card_top, card_width, card_height, value, label)

    # Charts
    imp_by_type = x_df.groupby("content_type")["impressions"].mean().sort_values(ascending=True).to_dict()
    eng_by_type = x_df.groupby("content_type")["engagement_rate"].mean().sort_values(ascending=True).to_dict()

    imp_chart = os.path.join(chart_dir, "x_imp_by_type.png")
    _create_bar_chart(imp_by_type, "Avg Impressions by Content Type", BLUE_HEX, imp_chart)
    chart_files.append(imp_chart)

    eng_chart = os.path.join(chart_dir, "x_eng_by_type.png")
    _create_bar_chart(eng_by_type, "Avg Engagement Rate by Content Type", BLUE_HEX, eng_chart)
    chart_files.append(eng_chart)

    slide.shapes.add_picture(imp_chart, Inches(0.5), Inches(3.2), Inches(6), Inches(4))
    slide.shapes.add_picture(eng_chart, Inches(6.8), Inches(3.2), Inches(6), Inches(4))

    return chart_files


def _build_x_top_posts_slide(prs: Presentation, x_df: pd.DataFrame) -> None:
    """Build Slide 6: X Top Posts table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.7),
                 "Top Performing X Posts", font_size=28, bold=True, color=BLUE)

    columns = [
        ("date", "Date", 1.8),
        ("caption_preview", "Tweet Preview", 5.5),
        ("impressions", "Impressions", 2.0),
        ("engagement_rate", "Eng. Rate", 2.0),
    ]
    _build_top_posts_slide(slide, x_df, columns)


def _build_insights_slide(
    prs: Presentation,
    ig_df: Optional[pd.DataFrame],
    x_df: Optional[pd.DataFrame],
    am_notes: Optional[dict] = None,
) -> None:
    """Build Slide 7: Key Insights with auto-generated bullets and AM notes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.7),
                 "Key Insights", font_size=28, bold=True, color=BLUE)

    insights = []

    if ig_df is not None and len(ig_df) > 0:
        # Top content type by engagement
        type_eng = ig_df.groupby("content_type")["engagement_rate"].mean()
        best_type = type_eng.idxmax()
        best_rate = round(type_eng.max(), 1)
        insights.append(
            f"{best_type}s drove the highest average engagement at {best_rate}%"
        )

        # Best day of week
        ig_copy = ig_df.copy()
        ig_copy["day_of_week"] = pd.to_datetime(ig_copy["date"]).dt.day_name()
        day_reach = ig_copy.groupby("day_of_week")["reach"].mean()
        best_day = day_reach.idxmax()
        best_day_reach = int(day_reach.max())
        insights.append(
            f"{best_day} was the strongest posting day with {best_day_reach:,} average reach"
        )

        # Top post
        top_post = ig_df.loc[ig_df["engagement_rate"].idxmax()]
        insights.append(
            f"The top post reached {int(top_post['reach']):,} people "
            f"with {top_post['engagement_rate']}% engagement"
        )

        # Platform summary
        insights.append(
            f"Instagram generated {int(ig_df['reach'].sum()):,} total reach "
            f"across {len(ig_df)} posts"
        )

    if x_df is not None and len(x_df) > 0:
        insights.append(
            f"X generated {int(x_df['impressions'].sum()):,} impressions "
            f"across {len(x_df)} tweets"
        )

    # Build bullet list
    txBox = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, insight in enumerate(insights):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = insight
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.font.name = FONT_NAME
        p.space_after = Pt(14)
        p.level = 0
        pf = p._pPr
        if pf is None:
            from pptx.oxml.ns import qn
            from lxml import etree
            pf = etree.SubElement(p._p, qn("a:pPr"))
        # Add bullet
        from pptx.oxml.ns import qn
        from lxml import etree
        buChar = etree.SubElement(pf, qn("a:buChar"))
        buChar.set("char", "\u2022")

    # AM Notes block (if provided and non-empty)
    if am_notes:
        sentiment = am_notes.get("sentiment", "")
        wins = am_notes.get("wins", "").strip()
        losses = am_notes.get("losses", "").strip()
        risks = am_notes.get("risks", "").strip()

        notes_lines = []
        if sentiment:
            notes_lines.append(f"Client Sentiment: {sentiment}")
        if wins:
            notes_lines.append(f"Wins: {wins}")
        if losses:
            notes_lines.append(f"Focus Areas: {losses}")
        if risks:
            notes_lines.append(f"Risks: {risks}")

        if notes_lines:
            _add_textbox(slide, Inches(0.9), Inches(5.8), Inches(11), Inches(0.35),
                         "AM Notes", font_size=13, bold=True, color=BLUE)
            _add_textbox(slide, Inches(0.9), Inches(6.15), Inches(11), Inches(1.1),
                         "  |  ".join(notes_lines),
                         font_size=11, color=GRAY)
    else:
        _add_textbox(slide, Inches(0.9), Inches(6.0), Inches(11), Inches(0.6),
                     "No AM notes provided for this period.",
                     font_size=12, color=GRAY, alignment=PP_ALIGN.LEFT)


def _add_health_badge(
    slide,
    left: int,
    top: int,
    score: int,
    band: str,
    color_hex: str,
) -> None:
    """Add a colored health score badge (rounded rect with score + band label).

    Args:
        slide: python-pptx slide object.
        left: Left position.
        top: Top position.
        score: Health score 0-100.
        band: Band label (Thriving / Healthy / At Risk / Critical).
        color_hex: Hex color string like "#27AE60".
    """
    from pptx.enum.shapes import MSO_SHAPE

    r = int(color_hex[1:3], 16)
    g = int(color_hex[3:5], 16)
    b = int(color_hex[5:7], 16)
    badge_color = RGBColor(r, g, b)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.9), Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = badge_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = False
    tf.auto_size = None

    p_score = tf.paragraphs[0]
    p_score.text = str(score)
    p_score.font.size = Pt(30)
    p_score.font.bold = True
    p_score.font.color.rgb = WHITE
    p_score.font.name = FONT_NAME
    p_score.alignment = PP_ALIGN.CENTER

    p_band = tf.add_paragraph()
    p_band.text = band
    p_band.font.size = Pt(11)
    p_band.font.color.rgb = WHITE
    p_band.font.name = FONT_NAME
    p_band.alignment = PP_ALIGN.CENTER


def _build_weekly_snapshot_slide(
    prs: Presentation,
    ig_df: Optional[pd.DataFrame],
    x_df: Optional[pd.DataFrame],
    deltas: dict,
    health: dict,
) -> None:
    """Build Weekly Slide 2: This Week at a Glance.

    Scorecards with WoW deltas + health score badge + top post callout.

    Args:
        prs: Presentation object.
        ig_df: IG DataFrame (full period).
        x_df: X DataFrame (full period).
        deltas: Delta dict from analytics.compute_deltas().
        health: Health score dict from analytics.compute_health_score().
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.7), Inches(0.3), Inches(9), Inches(0.65),
                 "This Week at a Glance", font_size=26, bold=True, color=BLUE)

    # Health badge top-right
    _add_health_badge(slide, Inches(11.1), Inches(0.25),
                      health["score"], health["band"], health["color"])

    def _delta_str(key: str) -> str:
        d = deltas.get(key, {})
        pct = d.get("delta_pct")
        if pct is None:
            return ""
        arrow = "+" if pct > 0 else ""
        return f"  {arrow}{pct}% WoW"

    # Scorecard row
    card_configs = []
    if ig_df is not None and len(ig_df) > 0:
        card_configs.extend([
            (f"{len(ig_df)}", "IG Posts" + _delta_str("ig_posts")),
            (f"{int(ig_df['reach'].sum()):,}", "IG Reach" + _delta_str("ig_reach")),
            (f"{round(ig_df['engagement_rate'].mean(), 2)}%", "IG Avg Eng. Rate" + _delta_str("ig_engagement_rate")),
        ])
    if x_df is not None and len(x_df) > 0:
        card_configs.extend([
            (f"{len(x_df)}", "Tweets" + _delta_str("x_posts")),
            (f"{int(x_df['impressions'].sum()):,}", "Impressions" + _delta_str("x_impressions")),
        ])

    num = min(len(card_configs), 4)
    if num > 0:
        card_w = Inches(2.9)
        card_h = Inches(1.5)
        card_top = Inches(1.2)
        gap = Inches(0.3)
        total_w = num * card_w + (num - 1) * gap
        start = (SLIDE_WIDTH - total_w) / 2
        for i, (val, lbl) in enumerate(card_configs[:num]):
            _add_scorecard(slide, start + i * (card_w + gap), card_top, card_w, card_h, val, lbl)

    # Top post callout box (IG only)
    if ig_df is not None and len(ig_df) > 0:
        top = ig_df.loc[ig_df["engagement_rate"].idxmax()]
        _add_textbox(slide, Inches(0.7), Inches(3.1), Inches(11.5), Inches(0.4),
                     "Top Post This Week", font_size=14, bold=True, color=BLUE)

        from pptx.enum.shapes import MSO_SHAPE
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(3.6), Inches(11.5), Inches(1.7),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BEIGE
        box.line.color.rgb = BLUE
        box.line.width = Pt(1)

        caption = str(top.get("caption_preview", ""))[:120]
        _add_textbox(slide, Inches(0.95), Inches(3.75), Inches(11), Inches(0.45),
                     f"“{caption}”", font_size=12, color=DARK_TEXT)
        meta = (
            f"{top['content_type']}  •  "
            f"Reach: {int(top['reach']):,}  •  "
            f"Engagement: {top['engagement_rate']}%  •  "
            f"{top['date']}"
        )
        _add_textbox(slide, Inches(0.95), Inches(4.45), Inches(11), Inches(0.4),
                     meta, font_size=11, color=GRAY)

    # Deliverables footer bar
    _add_textbox(slide, Inches(0.7), Inches(5.7), Inches(11.5), Inches(0.35),
                 "Powered by Ascnd ReportEngine",
                 font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)


def _build_weekly_content_slide(
    prs: Presentation,
    ig_df: Optional[pd.DataFrame],
    chart_dir: str,
) -> List[str]:
    """Build Weekly Slide 3: Content Performance breakdown.

    Args:
        prs: Presentation object.
        ig_df: IG DataFrame.
        chart_dir: Temp dir for chart PNGs.

    Returns:
        List of chart file paths created.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_files = []

    _add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11), Inches(0.65),
                 "Content Performance", font_size=26, bold=True, color=BLUE)

    if ig_df is None or len(ig_df) == 0:
        _add_textbox(slide, Inches(0.7), Inches(2.0), Inches(11), Inches(1.0),
                     "No Instagram data for this period.", font_size=16, color=GRAY)
        return chart_files

    reach_by_type = ig_df.groupby("content_type")["reach"].mean().sort_values(ascending=True).to_dict()
    eng_by_type = ig_df.groupby("content_type")["engagement_rate"].mean().sort_values(ascending=True).to_dict()

    reach_chart = os.path.join(chart_dir, "w_ig_reach.png")
    eng_chart = os.path.join(chart_dir, "w_ig_eng.png")
    _create_bar_chart(reach_by_type, "Avg Reach by Content Type", BLUE_HEX, reach_chart)
    _create_bar_chart(eng_by_type, "Avg Engagement Rate by Content Type", BLUE_HEX, eng_chart)
    chart_files.extend([reach_chart, eng_chart])

    slide.shapes.add_picture(reach_chart, Inches(0.5), Inches(1.2), Inches(6.0), Inches(3.8))
    slide.shapes.add_picture(eng_chart, Inches(6.9), Inches(1.2), Inches(6.0), Inches(3.8))

    # Compact breakdown table
    breakdown = (
        ig_df.groupby("content_type")
        .agg(posts=("post_id", "count"), avg_reach=("reach", "mean"), avg_eng=("engagement_rate", "mean"))
        .round(1)
        .sort_values("avg_eng", ascending=False)
        .reset_index()
    )
    cols_def = [
        ("content_type", "Type", 2.0),
        ("posts", "Posts", 1.2),
        ("avg_reach", "Avg Reach", 2.2),
        ("avg_eng", "Avg Eng. Rate", 2.2),
    ]
    rows = len(breakdown) + 1
    table_shape = slide.shapes.add_table(rows, 4, Inches(2.2), Inches(5.25), Inches(8.8), Inches(0.45 * rows))
    table = table_shape.table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(1.2)
    table.columns[2].width = Inches(2.8)
    table.columns[3].width = Inches(2.8)

    headers = ["Content Type", "Posts", "Avg Reach", "Avg Eng. Rate"]
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for ri, row in breakdown.iterrows():
        vals = [row["content_type"], str(int(row["posts"])),
                f"{int(row['avg_reach']):,}", f"{row['avg_eng']}%"]
        for ci, v in enumerate(vals):
            cell = table.cell(ri + 1, ci)
            cell.text = v
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_TEXT
            p.font.name = FONT_NAME
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BEIGE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return chart_files


def _build_weekly_am_slide(
    prs: Presentation,
    am_notes: Optional[dict],
    health: dict,
    goals_actuals: dict,
) -> None:
    """Build Weekly Slide 4: Account Summary with AM notes and health score.

    Args:
        prs: Presentation object.
        am_notes: AM qualitative inputs dict.
        health: Health score dict from analytics.compute_health_score().
        goals_actuals: Goals vs actuals dict.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.7), Inches(0.3), Inches(9), Inches(0.65),
                 "Account Summary", font_size=26, bold=True, color=BLUE)

    _add_health_badge(slide, Inches(11.1), Inches(0.25),
                      health["score"], health["band"], health["color"])

    notes = am_notes or {}
    sentiment = notes.get("sentiment", "")
    wins = notes.get("wins", "").strip()
    losses = notes.get("losses", "").strip()
    risks = notes.get("risks", "").strip()
    client_feedback = notes.get("client_feedback", "").strip()

    # Sentiment pill (colored label)
    sentiment_colors = {
        "Positive": BLUE_HEX,
        "Neutral": "#7F8C8D",
        "Negative": "#E67E22",
        "At Risk": "#E74C3C",
    }
    if sentiment:
        s_color_hex = sentiment_colors.get(sentiment, "#7F8C8D")
        from pptx.enum.shapes import MSO_SHAPE
        r_s = int(s_color_hex[1:3], 16)
        g_s = int(s_color_hex[3:5], 16)
        b_s = int(s_color_hex[5:7], 16)
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(0.7), Inches(1.1), Inches(2.4), Inches(0.5))
        pill.fill.solid()
        pill.fill.fore_color.rgb = RGBColor(r_s, g_s, b_s)
        pill.line.fill.background()
        tf_p = pill.text_frame
        tf_p.paragraphs[0].text = f"Sentiment: {sentiment}"
        tf_p.paragraphs[0].font.size = Pt(12)
        tf_p.paragraphs[0].font.bold = True
        tf_p.paragraphs[0].font.color.rgb = WHITE
        tf_p.paragraphs[0].font.name = FONT_NAME
        tf_p.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Two-column notes layout
    col_w = Inches(5.5)
    col_gap = Inches(0.5)

    def _notes_block(title: str, body: str, left: int, top: int) -> None:
        if not body:
            return
        _add_textbox(slide, left, top, col_w, Inches(0.35),
                     title, font_size=12, bold=True, color=BLUE)
        _add_textbox(slide, left, top + Inches(0.4), col_w, Inches(1.3),
                     body, font_size=11, color=DARK_TEXT)

    _notes_block("Strategic Wins", wins, Inches(0.7), Inches(1.75))
    _notes_block("Upcoming Risks", risks, Inches(0.7), Inches(3.4))
    _notes_block("Focus Areas", losses, Inches(6.9), Inches(1.75))
    _notes_block("Client Feedback", client_feedback, Inches(6.9), Inches(3.4))

    # Goals vs actuals footer
    if goals_actuals:
        _add_textbox(slide, Inches(0.7), Inches(5.5), Inches(11.5), Inches(0.35),
                     "Goals vs Actuals", font_size=12, bold=True, color=BLUE)
        goal_items = [f"{lbl}: {d['actual']:,} / {d['goal']:,} ({d['pct']}%{'  ' + chr(10003) if d['met'] else ''})"
                      for lbl, d in goals_actuals.items()]
        _add_textbox(slide, Inches(0.7), Inches(5.9), Inches(11.5), Inches(0.6),
                     "   |   ".join(goal_items), font_size=11, color=GRAY)

    _add_textbox(slide, Inches(0.7), Inches(6.8), Inches(11.5), Inches(0.3),
                 "Prepared by TNT | Powered by Ascnd",
                 font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)


def _build_closing_slide(prs: Presentation, date_range: str) -> None:
    """Build Slide 8: Closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BEIGE)

    _add_textbox(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2),
                 "Thank You", font_size=32, bold=True, color=BLUE,
                 alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.6),
                 "Prepared by TNT", font_size=20, color=DARK_TEXT,
                 alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(1), Inches(4.3), Inches(11.3), Inches(0.5),
                 "Report generated by Ascnd ReportEngine",
                 font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(1), Inches(5.1), Inches(11.3), Inches(0.5),
                 date_range, font_size=14, color=GRAY,
                 alignment=PP_ALIGN.CENTER)


def generate_report(
    ig_df: Optional[pd.DataFrame],
    x_df: Optional[pd.DataFrame],
    brand_name: str,
    output_path: str,
    report_type: str = "monthly",
    am_notes: Optional[dict] = None,
    deltas: Optional[dict] = None,
    goals_actuals: Optional[dict] = None,
    health: Optional[dict] = None,
) -> str:
    """Generate a branded PPTX report for a single brand.

    Dispatches to a 4-slide weekly format or an 8-slide monthly format
    based on report_type.

    Args:
        ig_df: Processed Instagram DataFrame, or None if no IG data.
        x_df: Processed X/Twitter DataFrame, or None if no X data.
        brand_name: Brand name for the report title.
        output_path: File path to save the PPTX file.
        report_type: "weekly" (4 slides) or "monthly" (8 slides).
        am_notes: AM qualitative inputs {sentiment, wins, losses, risks, etc.}.
        deltas: Period-over-period delta dict from analytics.compute_deltas().
        goals_actuals: Goals vs actuals dict from analytics.compute_goals_actuals().
        health: Health score dict from analytics.compute_health_score().

    Returns:
        The output file path where the PPTX was saved.
    """
    logger.info("Generating %s PPTX report for %s", report_type, brand_name)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    date_range = _format_date_range(ig_df, x_df)
    period_label = "Weekly" if report_type == "weekly" else "Monthly"
    chart_files: List[str] = []
    chart_dir = tempfile.mkdtemp(prefix="reportengine_charts_")

    # Fallback health dict so slides never crash if caller omits it
    _health = health or {"score": 0, "band": "N/A", "color": "#95A5A6"}

    try:
        if report_type == "weekly":
            # 4-slide compact weekly format
            _build_cover_slide(prs, brand_name, date_range, period_label)
            _build_weekly_snapshot_slide(prs, ig_df, x_df, deltas or {}, _health)
            files = _build_weekly_content_slide(prs, ig_df, chart_dir)
            chart_files.extend(files)
            _build_weekly_am_slide(prs, am_notes, _health, goals_actuals or {})

        else:
            # 8-slide monthly format
            _build_cover_slide(prs, brand_name, date_range, period_label)
            _build_summary_slide(prs, ig_df, x_df, deltas=deltas, goals_actuals=goals_actuals)

            if ig_df is not None and len(ig_df) > 0:
                files = _build_ig_overview_slide(prs, ig_df, chart_dir)
                chart_files.extend(files)
                _build_ig_top_posts_slide(prs, ig_df)

            if x_df is not None and len(x_df) > 0:
                files = _build_x_overview_slide(prs, x_df, chart_dir)
                chart_files.extend(files)
                _build_x_top_posts_slide(prs, x_df)

            _build_insights_slide(prs, ig_df, x_df, am_notes=am_notes)
            _build_closing_slide(prs, date_range)

        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        logger.info("PPTX saved to: %s", output_path)

    finally:
        for f in chart_files:
            try:
                os.remove(f)
            except OSError:
                pass
        try:
            os.rmdir(chart_dir)
        except OSError:
            pass

    return output_path
